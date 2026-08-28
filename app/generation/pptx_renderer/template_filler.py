# file: app/generation/pptx_renderer/template_filler.py
"""
Fills the DVT PPTX template from a cv_json produced by
cv_json_builder.build_cv_json_from_selection.

Hybrid approach:
  - Slide duplication (pagination) is raw XML surgery (see slide_duplication.py)
    -- python-pptx has no API for it, and it must happen BEFORE the file is
    opened with python-pptx.
  - Everything else (text replacement, bullets, photo removal) uses
    python-pptx's object model for readability. Bullet formatting isn't
    exposed by python-pptx's high-level API, so add_bullet_format() drops
    to paragraph._p (the underlying lxml element) for that one attribute --
    this is the documented, supported way to reach OOXML python-pptx doesn't
    wrap.

Shape IDs (SLIDE1_SHAPES / SLIDE2_SHAPES) were identified by inspecting the
template's raw slide XML -- specific to this template file, will break if
shapes are re-authored (id changes, shapes deleted/reordered).
"""

import shutil
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt
from sympy import exp

import ast

from .slide_duplication import duplicate_slide2

INVALID_VALUES = {
    "",
    "n/a",
    "na",
    "not specified",
    "not available",
    "unknown",
    'Unknown'
    "none",
    "-",
    "--",
}


YEARS_LABEL = {
    "French": "ans d'expérience",
    "English": "years of experience",
    "Spanish": "años de experiencia",
    "German": "Jahre Erfahrung",
}

PROJECT_ROOT = Path(__file__).resolve().parents[3]
TEMPLATE_PATH = PROJECT_ROOT / "Templates" / "Template_CV_format_DVT.pptx"

SLIDE1_SHAPES = {
    "name": 869, "summary": 870, "years": 872, "title": 877,
    "education": 879, "skills": 881, "exp1": 887, "exp2": 882,
    "photo": 2,
}
SLIDE2_SHAPES = {
    "name": 11,
    "title": 15,
    "years": 12,    
    "summary": 7,   
    "exp_box_a": 3,
    "exp_box_b": 5,
    "photo": 16,
}

CHAR_BUDGET_PER_BOX = 650
CHAR_BUDGET_PER_SLOT = 550
MAX_BULLETS_PER_EXPERIENCE = 5
MAX_CHARS_PER_BULLET = 220
SKILLS_CHAR_BUDGET = 500
SAFE_COLUMN_BOTTOM_EMU = 6300000

#clean text:
def clean_value(value):
    if value is None:
        return ""

    text = str(value).strip()

    if text.lower() in INVALID_VALUES:
        return ""

    return text 


# ---------------------------------------------------------------------------
# Shape lookup / low-level helpers
# ---------------------------------------------------------------------------

def find_shape_by_id(shapes, shape_id: int):
    """Recurses into groups. Returns None if not found on this slide."""
    
    for shape in shapes:
        if shape.shape_id == shape_id:
            return shape
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            found = find_shape_by_id(shape.shapes, shape_id)
            if found is not None:
                return found
    return None


def set_single_run_text(shape, text: str) -> None:
    """Replaces the text of a shape meant to hold ONE simple run, keeping
    the first run's formatting (font/size/bold) untouched."""
    text=clean_value(text)
    tf = shape.text_frame
    first_para = tf.paragraphs[0]
    if not first_para.runs:
        first_para.text = text or ""
        return
    first_para.runs[0].text = text or ""
    # Drop any extra runs/paragraphs so leftover template text can't survive
    # a shorter replacement value.
    for run in first_para.runs[1:]:
        run._r.getparent().remove(run._r)
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)


def add_bullet_format(paragraph, char: str = "•", indent_emu: int = 128588) -> None:
    """python-pptx has no high-level bullet API -- this reaches into the
    paragraph's underlying XML element (the documented escape hatch) to set
    marL/indent + buChar, matching the template's existing bullet style."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(indent_emu))
    pPr.set("indent", str(-indent_emu))
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(buFont)
    pPr.append(buChar)


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


# ---------------------------------------------------------------------------
# Experience formatting
# ---------------------------------------------------------------------------

def _experience_title_line(exp: dict) -> str:
    company = clean_value(exp.get("company"))
    title = clean_value(exp.get("title"))
    if company and title:
        return f"{company} | {title} :"
    return company or title or ""


def _experience_bullets(exp: dict) -> list[str]:
    print("EXP TYPE:", type(exp))
    print("EXP VALUE:", exp)
    bullets = []

    responsibilities = exp.get("responsibilities") or []

    if responsibilities:
        for r in responsibilities:
            print("RESP TYPE:", type(r))
            print("RESP VALUE:", r)
            if isinstance(r, dict):
                text = clean_value(
                    r.get("description") or r.get("category")
                )
            else:
                text = clean_value(r)

            if text:
                bullets.append(text)

    elif exp.get("description"):
        text = clean_value(exp.get("description"))
        if text:
            bullets.append(text)

    return bullets


def _experience_char_count(exp: dict) -> int:
    return len(_experience_title_line(exp)) + sum(len(b) for b in _experience_bullets(exp))


def _truncate_experience_for_display(exp: dict, char_budget: int) -> dict:
    exp = dict(exp)
    bullets = _experience_bullets(exp)[:MAX_BULLETS_PER_EXPERIENCE]
    bullets = [
        (b if len(b) <= MAX_CHARS_PER_BULLET else b[:MAX_CHARS_PER_BULLET].rstrip() + "…")
        for b in bullets
    ]
    exp["responsibilities"] = [{"description": b} for b in bullets]
    exp["description"] = ""
    return exp


def _paginate_by_char_budget(experiences: list[dict], budget_per_box: int) -> list[list[dict]]:
    boxes: list[list[dict]] = []
    current: list[dict] = []
    current_chars = 0
    for exp in experiences:
        exp_chars = _experience_char_count(exp)
        if current and current_chars + exp_chars > budget_per_box:
            boxes.append(current)
            current, current_chars = [], 0
        current.append(exp)
        current_chars += exp_chars
    if current:
        boxes.append(current)
    return boxes


def _write_experiences_into_shape(shape, experiences: list[dict], with_spacer: bool = False) -> None:
    if shape is None:
        return

    tf = shape.text_frame
    tf.clear()  # Efface le texte mais laisse 1 paragraphe vide à l'index 0

    # On récupère les bullets une seule fois par expérience
    first_exp = True

    for exp in experiences:
        bullets = _experience_bullets(exp)
        title_text = _experience_title_line(exp)

        # Si ce n'est pas la toute première expérience et qu'on veut un espacement
        if with_spacer and not first_exp:
            spacer_para = tf.add_paragraph()
            spacer_para.text = ""

        # Titre de l'expérience
        if first_exp:
            para = tf.paragraphs[0]
            first_exp = False
        else:
            para = tf.add_paragraph()

        if title_text:
            run = para.add_run()
            run.text = title_text
            run.font.bold = True
            run.font.size = Pt(8)

        # Puces de l'expérience
        for bullet_text in bullets:
            b_para = tf.add_paragraph()
            b_run = b_para.add_run()
            b_run.text = bullet_text
            b_run.font.size = Pt(8)
            add_bullet_format(b_para)


def fill_single_experience_shape(shape, exp: dict) -> None:
    _write_experiences_into_shape(shape, [exp] if exp else [], with_spacer=False)


def fill_multi_experience_shape(shape, experiences: list[dict]) -> None:
    _write_experiences_into_shape(shape, experiences, with_spacer=True)


# ---------------------------------------------------------------------------
# Static sections
# ---------------------------------------------------------------------------

def fill_skills_shape_dynamic(slide, shape_id: int, skills: list[str], font_pt: float = 8) -> None:
    """
    Fills the skills shape with ALL skills (no truncation), then grows the
    shape's height to fit the wrapped text and pushes every shape below it
    down by the same amount. Same approach as fill_education_shape_dynamic.
    """
    shape = find_shape_by_id(slide.shapes, shape_id)
    if shape is None:
        return

    tf = shape.text_frame
    tf.clear()
    clean = [str(s) for s in skills if s]
    if not clean:
        return

    text = "  •  ".join(clean)
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_pt)

    original_top = shape.top
    original_height = shape.height
    original_bottom = original_top + original_height

    needed_height = _estimate_needed_height_emu([text], shape.width, font_pt)
    new_height = max(original_height, needed_height)
    delta = new_height - original_height

    if delta <= 0:
        return

    shape.height = new_height

    for other in slide.shapes:
        if other.shape_id == shape_id:
            continue
        if other.top is None:
            continue
        if other.top >= original_bottom - 914400 // 20:
            other.top = other.top + delta


def fill_skills_shape(shape, skills: list[str]) -> None:
    """One flowing line ('Python  •  SQL  •  ...') instead of one bullet per
    skill. Truncated to SKILLS_CHAR_BUDGET -- past that, text overlaps the
    Formation section below it (observed empirically)."""
    tf = shape.text_frame
    tf.clear()
    clean = [str(s) for s in skills if s]

    kept = []
    running_len = 0
    separator_len = len("  •  ")
    for skill in clean:
        added_len = len(skill) + (separator_len if kept else 0)
        if running_len + added_len > SKILLS_CHAR_BUDGET:
            break
        kept.append(skill)
        running_len += added_len

    text = "  •  ".join(kept)
    

    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(8)

#Estimates the vertical space (EMU) needed to display text_lines

def _estimate_needed_height_emu(text_lines: list[str], shape_width_emu: int, font_pt: float) -> int:
    box_width_in = shape_width_emu / 914400
    avg_char_width_in = (0.5 * font_pt) / 72
    chars_per_line = max(1, int(box_width_in / avg_char_width_in))
    line_height_in = (1.2 * font_pt) / 72

    total_lines = 0
    for line in text_lines:
        total_lines += max(1, -(-len(line) // chars_per_line))  # ceil division -- wrapped line count

    return int(total_lines * line_height_in * 914400)

def fill_education_shape_dynamic(slide, shape_id: int, education: list[dict], font_pt: float = 7.5) -> None:
    shape = find_shape_by_id(slide.shapes, shape_id)
    if shape is None:
        return

    tf = shape.text_frame
    tf.clear()
    if not education:
        return

    lines = []
    first = True
    for edu in education:
        parts = [p for p in (edu.get("degree"), edu.get("field_of_study"), edu.get("institution")) if p]
        line = " - ".join(parts)
        if edu.get("years"):
            line = f"{line} ({edu['years']})" if line else str(edu["years"])
        if not line:
            continue
        lines.append(line)
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_pt)

    if not lines:
        return

    original_top = shape.top
    original_left = shape.left
    original_height = shape.height
    original_bottom = original_top + original_height

    # Debug: dump every shape's left/top so we can see exactly why a given
    # shape does or doesn't qualify as "same column, below education".
    for s in slide.shapes:
        print(f"[DEBUG all-shapes] id={s.shape_id} left={s.left} top={s.top}")

    needed_height = _estimate_needed_height_emu(lines, shape.width, font_pt)
    new_height = max(original_height, needed_height)
    delta = new_height - original_height
    if delta <= 0:
        return

    shape.height = new_height

    column_tolerance = 914400 // 4  # ~0.25in -- same-column horizontal tolerance
    vertical_tolerance = 914400 // 20  # ~0.05in

    pushed = 0
    for other in slide.shapes:
        if other.shape_id == shape_id:
            continue
        if other.top is None or other.left is None:
            continue
        same_column = abs(other.left - original_left) <= column_tolerance
        below_education = other.top >= original_bottom - vertical_tolerance
        if same_column and below_education:
            other.top = other.top + delta
            pushed += 1
    print(f"[DEBUG education] pushed {pushed} shape(s) down by {delta} (column-filtered)")


def fill_education_shape(shape, education: list[dict], font_pt: float = 7.5) -> None:
    capacity = _estimate_capacity_chars(shape, font_pt)
    tf = shape.text_frame
    tf.clear()
    if not education:
        return

    first = True
    used = 0
    for edu in education:
        parts = [p for p in (edu.get("degree"), edu.get("field_of_study"), edu.get("institution")) if p]
        line = " - ".join(parts)
        if edu.get("years"):
            line = f"{line} ({edu['years']})" if line else str(edu["years"])
        if not line:
            continue

        # Stop adding entries once the box's estimated capacity is exhausted --
        # prevents a long education list from overflowing into the section below.
        if used + len(line) > capacity and not first:
            break

        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_pt)
        used += len(line)


# ---------------------------------------------------------------------------
# Slide-level fill
# ---------------------------------------------------------------------------


def _fill_left_column_dynamic(slide, cv_json: dict) -> None:
    """
    Rebuilds the entire left column's vertical flow from scratch, instead of
    trusting the template's original box heights/positions (which are
    oversized placeholders relying on short text -- Formation's label (874)
    sits ABOVE skills' (881) own declared bottom edge even in the untouched
    template, so "push whatever is below original_bottom" is not a valid
    strategy here).

    Order (top to bottom, confirmed from template geometry): Compétences
    label (883) -> skills content (881) -> Formation label (874) ->
    education content (879) -> Réalisations label (885) -> exp2 block (887).
    """
    shapes = slide.shapes
    label_skills = find_shape_by_id(shapes, 883)
    skills_shape = find_shape_by_id(shapes, SLIDE1_SHAPES["skills"])
    label_formation = find_shape_by_id(shapes, 874)
    education_shape = find_shape_by_id(shapes, SLIDE1_SHAPES["education"])
    label_realisations = find_shape_by_id(shapes, 885)
    exp2_shape = find_shape_by_id(shapes, SLIDE1_SHAPES["exp2"])
    exp1_shape = find_shape_by_id(shapes, SLIDE1_SHAPES["exp1"])

    gap = 45720  # ~0.05in between stacked blocks

    # --- Fill skills content, compute its real needed height ---
    skills_text = ""
    if skills_shape is not None:
        skills = cv_json.get("skills") or []
        tf = skills_shape.text_frame
        tf.clear()
        clean = [str(s) for s in skills if s]
        if clean:
            skills_text = "  •  ".join(clean)
            run = tf.paragraphs[0].add_run()
            run.text = skills_text
            run.font.size = Pt(8)
    skills_height = (
        max(skills_shape.height, _estimate_needed_height_emu([skills_text], skills_shape.width, 8))
        if skills_shape is not None and skills_text else
        (skills_shape.height if skills_shape is not None else 0)
    )

    # --- Fill education content, compute its real needed height ---
    education_lines = []
    if education_shape is not None:
        education = cv_json.get("education") or []
        tf = education_shape.text_frame
        tf.clear()
        first = True
        for edu in education:
            parts = [p for p in (edu.get("degree"), edu.get("field_of_study"), edu.get("institution")) if p]
            line = " - ".join(parts)
            if edu.get("years"):
                line = f"{line} ({edu['years']})" if line else str(edu["years"])
            if not line:
                continue
            education_lines.append(line)
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = para.add_run()
            run.text = line
            run.font.size = Pt(7.5)
    education_height = (
        max(education_shape.height, _estimate_needed_height_emu(education_lines, education_shape.width, 7.5))
        if education_shape is not None and education_lines else
        (education_shape.height if education_shape is not None else 0)
    )

    # --- Reflow: place each block right after the previous one's real height ---
    cursor = label_skills.top if label_skills is not None else 1167879

    if label_skills is not None:
        cursor = label_skills.top + label_skills.height + gap
    if skills_shape is not None:
        skills_shape.top = cursor
        cursor = cursor + skills_height + gap

    if label_formation is not None:
        label_formation.top = cursor
        cursor = cursor + label_formation.height + gap
    if education_shape is not None:
        education_shape.top = cursor
        cursor = cursor + education_height + gap

    if label_realisations is not None:
        label_realisations.top = cursor
        cursor = cursor + label_realisations.height + gap
    if exp2_shape is not None:
        exp2_shape.top = cursor

    # --- Relocate exp2 to the right column if the left column ran out of room ---
    if exp2_shape is not None and exp2_shape.top >= SAFE_COLUMN_BOTTOM_EMU:
        print(f"[layout] exp2 top={exp2_shape.top} exceeds safe bottom -- relocating to second column")
        if exp1_shape is not None:
            exp2_shape.left = exp1_shape.left
            exp2_shape.top = exp1_shape.top + exp1_shape.height + 91440


def fill_slide1(slide, cv_json: dict, target_language: str = "French") -> None:
    shapes = slide.shapes
    set_single_run_text(find_shape_by_id(shapes, SLIDE1_SHAPES["name"]), cv_json.get("name") or "")
    set_single_run_text(find_shape_by_id(shapes, SLIDE1_SHAPES["title"]), cv_json.get("title") or "")

    years = cv_json.get("years_of_experience")
    label = YEARS_LABEL.get(target_language, YEARS_LABEL["French"])
    set_single_run_text(
        find_shape_by_id(shapes, SLIDE1_SHAPES["years"]),
        f"{years} {label}" if years else "",
    )
    set_single_run_text(find_shape_by_id(shapes, SLIDE1_SHAPES["summary"]), cv_json.get("summary") or "")

    # 1. Remplir D'ABORD les contenus des expériences
    experiences = cv_json.get("experience") or []
    exp1_shape = find_shape_by_id(shapes, SLIDE1_SHAPES["exp1"])
    exp2_shape = find_shape_by_id(shapes, SLIDE1_SHAPES["exp2"])
    
    exp1 = experiences[0] if len(experiences) > 0 else {}
    exp2 = experiences[1] if len(experiences) > 1 else {}
    
    fill_single_experience_shape(exp1_shape, exp1)
    fill_single_experience_shape(exp2_shape, exp2)

    # 2. Repositionner dynamiquement la colonne de gauche (skills, edu, exp2) APRÈS le remplissage
    _fill_left_column_dynamic(slide, cv_json)

    photo = find_shape_by_id(shapes, SLIDE1_SHAPES["photo"])
    if photo is not None:
        remove_shape(photo)
   
    



def _get_slide2_box_capacities(template_path) -> tuple[int, int]:
    """
    Reads exp_box_a/exp_box_b's real dimensions from the template once,
    to replace the flat CHAR_BUDGET_PER_BOX constant. All duplicated slide-2
    pages share the same geometry, so this only needs to run once per render.
    """
    prs = Presentation(str(template_path))
    slide2 = prs.slides[1]
    box_a = find_shape_by_id(slide2.shapes, SLIDE2_SHAPES["exp_box_a"])
    box_b = find_shape_by_id(slide2.shapes, SLIDE2_SHAPES["exp_box_b"])
    return (
        _estimate_capacity_chars(box_a, font_pt=8),
        _estimate_capacity_chars(box_b, font_pt=8),
    )


def fill_slide2_page(slide, cv_json: dict, boxes: list[list[dict]], target_language: str = "French") -> None:
    shapes = slide.shapes
    
    # 1. Nom & Titre
    set_single_run_text(find_shape_by_id(shapes, SLIDE2_SHAPES["name"]), cv_json.get("name") or "")
    set_single_run_text(find_shape_by_id(shapes, SLIDE2_SHAPES["title"]), cv_json.get("title") or "")

    # 2. Années d'expérience
    years = cv_json.get("years_of_experience")
    label = YEARS_LABEL.get(target_language, YEARS_LABEL["French"])
    years_shape = find_shape_by_id(shapes, SLIDE2_SHAPES["years"])
    if years_shape:
        set_single_run_text(years_shape, f"{years} {label}" if years else "")

    # 3. Résumé / Summary
    summary_shape = find_shape_by_id(shapes, SLIDE2_SHAPES["summary"])
    if summary_shape:
        set_single_run_text(summary_shape, cv_json.get("summary") or "")

    # 4. Expériences & Photo
    box_a = find_shape_by_id(shapes, SLIDE2_SHAPES["exp_box_a"])
    box_b = find_shape_by_id(shapes, SLIDE2_SHAPES["exp_box_b"])
    fill_multi_experience_shape(box_a, boxes[0] if len(boxes) > 0 else [])
    fill_multi_experience_shape(box_b, boxes[1] if len(boxes) > 1 else [])

    photo = find_shape_by_id(shapes, SLIDE2_SHAPES["photo"])
    if photo is not None:
        remove_shape(photo)


#Remove Extra Slide
def _remove_slide(prs, index: int) -> None:    
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def render_cv_pptx(cv_json: dict, output_path: str, target_language: str = "French") -> str:
    import uuid
    workdir = Path("/tmp") / f"pptx_render_{uuid.uuid4().hex}"
    box_a_capacity, box_b_capacity = _get_slide2_box_capacities(TEMPLATE_PATH)
    per_box_budget = min(box_a_capacity, box_b_capacity)
    if workdir.exists():
        shutil.rmtree(workdir)
    workdir.mkdir(parents=True)

    with zipfile.ZipFile(TEMPLATE_PATH) as z:
        z.extractall(workdir)

    # 1. Compute pagination BEFORE touching python-pptx.
    experiences = cv_json.get("experience") or []
    remaining = experiences[2:]
    boxes = _paginate_by_char_budget(remaining, per_box_budget)
    # No "or [[]]" fallback -- an empty pages list means slide 2 isn't needed at all.
    pages = [boxes[i:i + 2] for i in range(0, len(boxes), 2)]

    extra_slide_paths = [duplicate_slide2(workdir) for _ in pages[1:]]

    # 3. Rezip into an intermediate file, then open it with python-pptx.
    intermediate_path = workdir.parent / f"{workdir.name}_intermediate.pptx"
    with zipfile.ZipFile(intermediate_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for file in workdir.rglob("*"):
            if file.is_file():
                zf.write(file, file.relative_to(workdir))

    prs = Presentation(str(intermediate_path))

    # slides[0] = slide1, slides[1] = slide2, slides[2:] = duplicated pages,
    # in the same order they were added to sldIdLst.
    fill_slide1(prs.slides[0], cv_json, target_language)

    if pages:
        fill_slide2_page(prs.slides[1], cv_json, pages[0], target_language=target_language)
        for i, page in enumerate(pages[1:], start=2):
            fill_slide2_page(prs.slides[i], cv_json, page, target_language=target_language)
    else:
        _remove_slide(prs, index=1)

    prs.save(output_path)

    shutil.rmtree(workdir)
    intermediate_path.unlink()
    return output_path

#Rough estimate of how many characters fit in a shape's text box
def _estimate_capacity_chars(shape, font_pt: float = 8) -> int:   
    if shape is None or not shape.width or not shape.height:
        return 500  # fallback if shape geometry is unreadable

    box_width_in = shape.width / 914400
    box_height_in = shape.height / 914400

    avg_char_width_in = (0.5 * font_pt) / 72
    line_height_in = (1.2 * font_pt) / 72

    chars_per_line = max(1, box_width_in / avg_char_width_in)
    num_lines = max(1, box_height_in / line_height_in)

    return int(chars_per_line * num_lines)

